from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_FILL
from pptx.util import Pt
import lzma
import os

class EmbedExtract():
	"""A class that contains methods to embed and extract steganograms into a given cover PowerPoint file."""
	def __init__(self, pptx_file_name):
		super().__init__()
		self._pptx_file_name = pptx_file_name
		self._pptx_file = Presentation(self.pptx_file_name)
		self._statinfo = os.stat(self.pptx_file_name)
		self._st_atime_ns = self.statinfo.st_atime_ns
		self._st_mtime_ns = self.statinfo.st_mtime_ns
		self._ori_st_size = self.statinfo.st_size

	@property
	def pptx_file_name(self):
		return self._pptx_file_name

	@property
	def pptx_file(self):
		return self._pptx_file

	@property
	def statinfo(self):
		return self._statinfo

	@property
	def st_atime_ns(self):
		return self._st_atime_ns

	@property
	def st_mtime_ns(self):
		return self._st_mtime_ns

	@property
	def ori_st_size(self):
		return self._ori_st_size
	
	
	def save_pptx(self, file_name = None):
		if file_name is None:
			output_path = 'output\\' + os.path.basename(self.pptx_file_name)
			self.pptx_file.save(output_path)
			os.utime(output_path, ns=(self.st_atime_ns, self.st_mtime_ns))
		else:
			self.pptx_file.save(file_name)
		

	def embed_hex(self, hex_strings):

		if (self.calculate_capacity()) > len(hex_strings):
			hex_strings.reverse()
			pptx = self.pptx_file
			slides = pptx.slides

			for slide in slides:
				shapes = slide.shapes
				for shape in shapes:

					left = shape.left
					top = shape.top
					width = shape.width
					height = shape.height

					is_auto_shape = self.is_auto_shape(shape)
					is_as_outl_tr = False
					is_dim_usable = True
					has_text_frame = shape.has_text_frame

					if is_auto_shape:
						is_as_outl_tr = self.is_as_outl_tr(shape)

					if is_as_outl_tr:
						if len(hex_strings) == 0:
							line_width = 256
							shape.line.width = line_width
							return
						else:
							line_width = hex_strings.pop()
							line_width = int(line_width, 16)
							shape.line.width = line_width
							

					if left is None:
						is_dim_usable = False

					elif top is None:
						is_dim_usable = False

					elif width is None:
						is_dim_usable = False

					elif height is None:
						is_dim_usable = False

					if is_dim_usable:

						if len(hex_strings) == 0:
							left = left // 1000
							left = left * 1000
							left = left + 256
							self.change_shape_dimensions(shape, dimensions = [left, top, width, height])
							return
						else:
							left = left // 1000
							left = left * 1000
							left = left + int(hex_strings.pop(), 16)
							self.change_shape_dimensions(shape, dimensions = [left, top, width, height])
						
						if len(hex_strings) == 0:
							top = top // 1000
							top = top * 1000
							top = top + 256
							self.change_shape_dimensions(shape, dimensions = [left, top, width, height])
							return
						else:
							top = top // 1000
							top = top * 1000
							top = top + int(hex_strings.pop(), 16)
							self.change_shape_dimensions(shape, dimensions = [left, top	, width, height])

						if len(hex_strings) == 0:
							width = width // 1000
							width = width * 1000
							width = width + 256
							self.change_shape_dimensions(shape, dimensions = [left, top, width, height])
							return
						else:
							width = width // 1000
							width = width * 1000
							width = width + int(hex_strings.pop(), 16)
							self.change_shape_dimensions(shape, dimensions = [left, top, width, height])
							
						if len(hex_strings) == 0:
							height = height // 1000
							height = height * 1000
							height = height + 256
							self.change_shape_dimensions(shape, dimensions = [left, top, width, height])
							return
						else:
							height = height // 1000
							height = height * 1000
							height = height + int(hex_strings.pop(), 16)
							self.change_shape_dimensions(shape, dimensions = [left, top, width, height])

					if has_text_frame:
						text_frame = shape.text_frame
						first_paragraph = text_frame.paragraphs[0]
						run = first_paragraph.add_run()
						font = run.font
						font.size = Pt(1)
						font.fill.background()

						if len(hex_strings) == 0:
							text = '256'
							run.text = text
							return
						elif len(hex_strings) >= 4000:
							text = [hex_strings.pop() for i in range(0,4000)] 
							text = ''.join(text)
							run.text = text
						elif len(hex_strings) < 4000:
							hex_strings_length = len(hex_strings)
							text = [hex_strings.pop() for i in range(0,hex_strings_length)]
							text = ''.join(text)
							run.text = text

					if len(hex_strings) == 0:
						shape.name = '256'
						return
					elif len(hex_strings) >= 4000:
						text = [hex_strings.pop() for i in range(0,4000)]
						text = ''.join(text)
						shape.name = text
					elif len(hex_strings) < 4000:
						hex_strings_length = len(hex_strings)
						text = [hex_strings.pop() for i in range(0, hex_strings_length)]
						text = ''.join(text)
						shape.name = text

		else:
			raise InsufficientCapacityError(self.calculate_capacity(), len(hex_strings) + 1)

	def extract_hex(self):

		hex_strings = []
		pptx = self.pptx_file
		slides =  pptx.slides

		for slide in slides:
			shapes = slide.shapes
			for shape in shapes:

				left = shape.left
				top = shape.top
				width = shape.width
				height = shape.height

				is_auto_shape = self.is_auto_shape(shape)
				is_as_outl_tr = False
				is_dim_usable = True
				has_text_frame = shape.has_text_frame

				if is_auto_shape:
					is_as_outl_tr = self.is_as_outl_tr(shape)

				if is_as_outl_tr:
					line_width = shape.line.width
					if line_width == 256:
						return hex_strings
					else:
						line_width = bytes([line_width]).hex()
						hex_strings.append(line_width)

				if left is None:
					is_dim_usable = False

				elif top is None:
					is_dim_usable = False

				elif width is None:
					is_dim_usable = False

				elif height is None:
					is_dim_usable = False

				if is_dim_usable:

					if left % 1000 == 256:
						return hex_strings
					else:
						left = left % 1000
						left = bytes([left]).hex()
						hex_strings.append(left)

					if top % 1000 == 256:
						return hex_strings
					else:
						top = top % 1000
						top = bytes([top]).hex()
						hex_strings.append(top)

					if width % 1000 == 256:
						return hex_strings
					else:
						width = width % 1000
						width = bytes([width]).hex()
						hex_strings.append(width)

					if height % 1000 == 256:
						return hex_strings
					else:
						height = height % 1000
						height = bytes([height]).hex()
						hex_strings.append(height)

				if has_text_frame:

					text_frame = shape.text_frame
					first_paragraph = text_frame.paragraphs[0]
					run = first_paragraph.runs[-1]
					text = run.text

					if text == '256':
						return hex_strings
					else:
						text = [text[i:i+2] for i in range(0, len(text), 2)]
						hex_strings += text

				text = shape.name
				if text == '256':
					return hex_strings
				else:
					text = [text[i:i+2] for i in range(0, len(text), 2)]
					hex_strings += text

	def change_shape_dimensions(self, shape, dimensions):
		"""Helper method to change the shape dimensions"""
		shape.left, shape.top, shape.width, shape.height = dimensions

	def get_mod_st_size(self):
		"""Helper method to get the size of the saved pptx_file"""
		self.save_pptx('AB12EF34.pptx')
		statinfo = os.stat('AB12EF34.pptx')
		mod_st_size = statinfo.st_size
		os.remove('AB12EF34.pptx')
		return mod_st_size

	def get_st_size_diff(self):
		"""Helper method to get the difference between the size of the original pptx_file and the saved pptx_file."""
		return self.ori_st_size - self.get_mod_st_size()

	def get_num_of_shapes(self):
		"""Helper method to get the number of shapes in the pptx_file"""
		pptx = self.pptx_file
		count = 0
		slides = pptx.slides
		for slide in slides:
			shapes = slide.shapes
			count += len(shapes)
		return count

	def get_num_of_auto_shapes(self):
		"""Helper method to get the number of auto shapes in the pptx_file"""
		pptx = self.pptx_file
		count = 0
		slides = pptx.slides
		for slide in slides:
			shapes = slide.shapes
			for shape in shapes:
				if self.is_auto_shape(shape):
					count += 1
		return count

	def get_num_of_tr_outl(self):
		"""Helper method to get the number of transparent auto shape outlines"""
		pptx = self.pptx_file
		count = 0
		slides = pptx.slides
		for slide in slides:
			shapes = slide.shapes
			for shape in shapes:
				if self.is_auto_shape(shape):
					if self.is_as_outl_tr(shape):
						count += 1
		return count

	def get_num_of_dimensions(self):
		"""Helper method to get the number of dimensions available for use"""
		count = 0
		pptx = self.pptx_file
		slides = pptx.slides
		for slide in slides:
			shapes = slide.shapes
			for shape in shapes:
				if shape.left is None:
					continue
				elif shape.top is None:
					continue
				elif shape.width is None:
					continue
				elif shape.height is None:
					continue
				count = count + 4
		return count

	def is_auto_shape(self, shape):
		"""Helper method to determine if a shape is an auto shape"""
		is_auto_shape = False
		if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
			is_auto_shape = True
		return is_auto_shape

	def is_as_outl_tr(self, shape):
		"""Helper method to determine if the auto shape's outline is transparent"""
		is_as_outl_tr = False
		try:
			_ = shape.line.fill.fore_color
		except TypeError as e:
			shape.line.fill.background()
			is_as_outl_tr = True
		else:
			if shape.line.fill.type == MSO_FILL.BACKGROUND:
				is_as_outl_tr = True
		return is_as_outl_tr

	def get_num_of_text_frames(self):
		"""Helper method to get the number of text frames in pptx_file"""
		pptx = self.pptx_file
		count = 0
		slides = pptx.slides
		for slide in slides:
			shapes = slide.shapes
			for shape in shapes:
				if shape.has_text_frame:
					count += 1
		return count

	def calculate_capacity(self):
		"""Helper method to calculate the steganographic capacity of pptx_file"""
		capacity = self.get_num_of_tr_outl() + self.get_num_of_dimensions() + ((self.get_st_size_diff() // 1024) * 4000)
		if capacity < 0:
			return 0
		return capacity

class Error(Exception):
	"""Base class for exceptions in this module."""
	pass

class InsufficientCapacityError(Error):
	"""Exception raised when there is insufficient steganographic capacity
	provided by the cover PowerPoint file to store the steganograms.
	"""
	def __init__(self, steganographic_capacity, required_capacity):
		super().__init__()
		self.steganographic_capacity = steganographic_capacity
		self.required_capacity = required_capacity
		
	def __str__(self):
		avail_steg_cap = 'Available Steganographic Capacity: {} bytes\n'.format(self.steganographic_capacity)
		req_steg_cap = 'Required Steganographic Capacity: {} bytes\n'.format(self.required_capacity)
		return avail_steg_cap + req_steg_cap


##notes slides