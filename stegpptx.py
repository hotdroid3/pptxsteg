from pptx import Presentation
import math


def openpowerpointfile():
    """Open a PowerPoint file specified by user and returns a PowerPoint object."""
    filename = raw_input('Only PowerPoint 2007 files or later can be used!\nEnter path to PowerPoint file:\n')
    prs = Presentation(filename)
    return prs


def savepowerpointfile(prs):
    """Save a PowerPoint file with a name specified by user
    :type prs: Presentation object
    """
    filename = raw_input('Enter name of PowerPoint file to be saved: \n')
    prs.save(filename)
    return None


def requeststegmessage():
    """Request the user for the steganographic message the user wishes to embed and returns a list containing each
    character of the message"""
    msg = raw_input('Enter message to be embedded in chosen PowerPoint file: ')
    msg = msg.strip()
    msg = msg.encode('hex')
    msg = list(msg)
    return msg


def checkcapacity(msglist, prs):
    """Check if there is sufficient space to embed steganographic message"""
    slidescol = prs.slides
    numofslides = slidescol.__len__()
    numofplaceholders = 0
    for x in range(0,numofslides):
        slide = slidescol.__getitem__(x)
        numofplaceholders = numofplaceholders + slide.shapes.__len__()
    if msglist.__len__() > numofplaceholders:
        return False
    else:
        return True


def embedhex(msglist, prs):
    """Embed the hex value of message into PowerPoint file"""
    msglist.reverse()
    listindex = msglist.__len__() - 1
    slidescol = prs.slides
    numofslides = slidescol.__len__()
    for x in range(0,numofslides):
        slide = slidescol.__getitem__(x)
        shapescol = slide.shapes
        numofplaceholders = shapescol.__len__()
        for y in range(0, numofplaceholders):
            if msglist.__len__() == 0:
                break
            else:
                shape = shapescol.__getitem__(y)
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height
                #rotation = shape.rotation
                if msglist[listindex] == '0':
                    if left % 2 == 0:
                        left = left + 1
                    if top % 2 == 0:
                        top = top + 1
                    if width % 2 == 0:
                        width = width + 1
                    if height % 2 == 0:
                        height = height + 1
                elif msglist[listindex] == '1':
                    if left % 2 == 0:
                        left = left + 1
                    if top % 2 == 0:
                        top = top + 1
                    if width % 2 == 0:
                        width = width + 1
                    if height % 2 == 1:
                        height = height + 1
                elif msglist[listindex] == '2':
                    if left % 2 == 0:
                        left = left + 1
                    if top % 2 == 0:
                        top = top + 1
                    if width % 2 == 1:
                        width = width + 1
                    if height % 2 == 0:
                        height = height + 1
                elif msglist[listindex] == '3':
                    if left % 2 == 0:
                        left = left + 1
                    if top % 2 == 0:
                        top = top + 1
                    if width % 2 == 1:
                        width = width + 1
                    if height % 2 == 1:
                        height = height + 1
                elif msglist[listindex] == '4':
                    if left % 2 == 0:
                        left = left + 1
                    if top % 2 == 1:
                        top = top + 1
                    if width % 2 == 0:
                        width = width + 1
                    if height % 2 == 0:
                        height = height + 1
                elif msglist[listindex] == '5':
                    if left % 2 == 0:
                        left = left + 1
                    if top % 2 == 1:
                        top = top + 1
                    if width % 2 == 0:
                        width = width + 1
                    if height % 2 == 1:
                        height = height + 1
                elif msglist[listindex] == '6':
                    if left % 2 == 0:
                        left = left + 1
                    if top % 2 == 1:
                        top = top + 1
                    if width % 2 == 1:
                        width = width + 1
                    if height % 2 == 0:
                        height = height + 1
                elif msglist[listindex] == '7':
                    if left % 2 == 0:
                        left = left + 1
                    if top % 2 == 1:
                        top = top + 1
                    if width % 2 == 1:
                        width = width + 1
                    if height % 2 == 1:
                        height = height + 1
                elif msglist[listindex] == '8':
                    if left % 2 == 1:
                        left = left + 1
                    if top % 2 == 0:
                        top = top + 1
                    if width % 2 == 0:
                        width = width + 1
                    if height % 2 == 0:
                        height = height + 1
                elif msglist[listindex] == '9':
                    if left % 2 == 1:
                        left = left + 1
                    if top % 2 == 0:
                        top = top + 1
                    if width % 2 == 0:
                        width = width + 1
                    if height % 2 == 1:
                        height = height + 1
                elif msglist[listindex] == 'A' or msglist[listindex] == 'a':
                    if left % 2 == 1:
                        left = left + 1
                    if top % 2 == 0:
                        top = top + 1
                    if width % 2 == 1:
                        width = width + 1
                    if height % 2 == 0:
                        height = height + 1
                elif msglist[listindex] == 'B' or msglist[listindex] == 'b':
                    if left % 2 == 1:
                        left = left + 1
                    if top % 2 == 0:
                        top = top + 1
                    if width % 2 == 1:
                        width = width + 1
                    if height % 2 == 1:
                        height = height + 1
                elif msglist[listindex] == 'C' or msglist[listindex] == 'c':
                    if left % 2 == 1:
                        left = left + 1
                    if top % 2 == 1:
                        top = top + 1
                    if width % 2 == 0:
                        width = width + 1
                    if height % 2 == 0:
                        height = height + 1
                elif msglist[listindex] == 'D' or msglist[listindex] == 'd':
                    if left % 2 == 1:
                        left = left + 1
                    if top % 2 == 1:
                        top = top + 1
                    if width % 2 == 0:
                        width = width + 1
                    if height % 2 == 1:
                        height = height + 1
                elif msglist[listindex] == 'E' or msglist[listindex] == 'e':
                    if left % 2 == 1:
                        left = left + 1
                    if top % 2 == 1:
                        top = top + 1
                    if width % 2 == 1:
                        width = width + 1
                    if height % 2 == 0:
                        height = height + 1
                elif msglist[listindex] == 'F' or msglist[listindex] == 'f':
                    if left % 2 == 1:
                        left = left + 1
                    if top % 2 == 1:
                        top = top + 1
                    if width % 2 == 1:
                        width = width + 1
                    if height % 2 == 1:
                        height = height + 1

                shape.left = left
                shape.top = top
                shape.width = width
                shape.height = height
                #shape.rotation = rotation + 0.01
                shape.name = 'e'
                listindex = listindex - 1
                msglist.pop()

    return None
def extracthex(prs):
    """Extract hex message from PowerPoint file specified by user and return the hex value"""
    msglist = []
    slidescol = prs.slides
    numofslides = slidescol.__len__()
    for x in range(0,numofslides):
        slide = slidescol.__getitem__(x)
        shapescol = slide.shapes
        numofplaceholders = shapescol.__len__()
        for y in range(0, numofplaceholders):
            shape = shapescol.__getitem__(y)
            name = shape.name
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height
            #rotation = shape.rotation
            #tmprotation = math.trunc(rotation)
            #if(rotation - tmprotation) != 0.01:
            #    break
            if name != 'e':
                break
            else:
                if left % 2 == 1:
                    if top % 2 == 1:
                        if width % 2 == 1:
                            if height % 2 == 1:
                                msglist.append('0')
                            else:
                                msglist.append('1')
                        else:
                            if height % 2 == 1:
                                msglist.append('2')
                            else:
                                msglist.append('3')
                    else:
                        if width % 2 == 1:
                            if height % 2 == 1:
                                msglist.append('4')
                            else:
                                msglist.append('5')
                        else:
                            if height % 2 == 1:
                                msglist.append('6')
                            else:
                                msglist.append('7')
                else:
                    if top % 2 == 1:
                        if width % 2 == 1:
                            if height % 2 == 1:
                                msglist.append('8')
                            else:
                                msglist.append('9')
                        else:
                            if height % 2 == 1:
                                msglist.append('A')
                            else:
                                msglist.append('B')
                    else:
                        if width % 2 == 1:
                            if height % 2 == 1:
                                msglist.append('C')
                            else:
                                msglist.append('D')
                        else:
                            if height % 2 == 1:
                                msglist.append('E')
                            else:
                                msglist.append('F')


    return msglist

flag = True
while(flag):
    userchoice = raw_input("\n\nPlease select your choice:\n1. Encode\n2. Decode\n3. Exit\n\n")
    if userchoice == '1':
        prs = openpowerpointfile()
        msg = requeststegmessage()
        capflag = checkcapacity(msg,prs)
        if capflag:
            embedhex(msg,prs)
            savepowerpointfile(prs)
        else:
            print 'PowerPoint file does not have enough steganographic capacity'
    if userchoice == '2':
        prs = openpowerpointfile()
        msglist = extracthex(prs)
        msg = ''.join(msglist)
        msg = msg.decode('hex')
        print '\n\nMessage: ', msg
    if userchoice == '3':
        flag = False
    userchoice = 0




